"""
Утиліти для конвертації офісних документів у HTML для попереднього перегляду
"""
import os
from io import BytesIO
from typing import Optional, Tuple


def convert_docx_to_html(file_path: str) -> Tuple[str, Optional[str]]:
    """
    Конвертує .docx файл у HTML для відображення в браузері
    
    Args:
        file_path: Шлях до .docx файлу
        
    Returns:
        Tuple[str, Optional[str]]: (HTML контент, повідомлення про помилку)
    """
    try:
        from docx import Document
        from docx.shared import Pt
        
        doc = Document(file_path)
        html_parts = []
        
        html_parts.append('<div class="document-page">')
        
        # Обробка параграфів
        for para in doc.paragraphs:
            if not para.text.strip():
                html_parts.append('<p>&nbsp;</p>') # Empty line
                continue
                
            # Визначення стилю
            style_class = ""
            if para.style.name.startswith('Heading'):
                level = para.style.name.replace('Heading ', '')
                if level.isdigit():
                    html_parts.append(f'<h{level}>{para.text}</h{level}>')
                else:
                    html_parts.append(f'<h3>{para.text}</h3>')
            else:
                # Обробка форматування тексту
                text_html = para.text
                html_parts.append(f'<p>{text_html}</p>')
        
        # Обробка таблиць
        for table in doc.tables:
            html_parts.append('<table class="table table-bordered table-striped mt-3 excel-table">')
            for i, row in enumerate(table.rows):
                html_parts.append('<tr>')
                for cell in row.cells:
                    tag = 'th' if i == 0 else 'td'
                    cell_text = cell.text
                    # Add onclick handler to table cells for formula bar
                    if tag == 'td':
                        html_parts.append(f'<{tag} data-formula="{cell_text}" onclick="showFormula(this)">{cell_text}</{tag}>')
                    else:
                        html_parts.append(f'<{tag}>{cell_text}</{tag}>')
                html_parts.append('</tr>')
            html_parts.append('</table>')
            
        html_parts.append('</div>')
        
        html_content = '\n'.join(html_parts)
        
        if not html_content.strip() or html_content == '<div class="document-page"></div>':
            html_content = '<div class="document-page"><p class="text-muted">Документ порожній або не містить текстового контенту.</p></div>'
        
        return html_content, None
        
    except Exception as e:
        error_msg = f"Помилка при читанні документа: {str(e)}"
        return "", error_msg


def convert_xlsx_to_html(file_path: str, max_rows: int = 100) -> Tuple[str, Optional[str]]:
    """
    Конвертує .xlsx файл у HTML таблиці з підтримкою формул та стилів
    """
    try:
        from openpyxl import load_workbook
        from openpyxl.styles import PatternFill, Font
        
        # Load twice: once for values, once for formulas
        wb_values = load_workbook(file_path, data_only=True)
        wb_formulas = load_workbook(file_path, data_only=False)
        
        html_parts = []
        
        for sheet_name in wb_values.sheetnames:
            sheet_values = wb_values[sheet_name]
            sheet_formulas = wb_formulas[sheet_name]
            
            # Wrap each sheet in a document page
            html_parts.append('<div class="document-page excel-sheet">')
            html_parts.append(f'<h4 class="sheet-title mb-3">{sheet_name}</h4>')
            
            # Find data boundaries
            max_row = min(sheet_values.max_row, max_rows)
            max_col = sheet_values.max_column
            
            if max_row == 0 or max_col == 0:
                html_parts.append('<p class="text-muted">Аркуш порожній</p>')
                html_parts.append('</div>')
                continue
            
            html_parts.append('<div class="table-responsive">')
            html_parts.append('<table class="excel-table">')
            
            # Column headers (A, B, C...)
            html_parts.append('<thead><tr><th class="row-header"></th>')
            for col in range(1, max_col + 1):
                from openpyxl.utils import get_column_letter
                col_letter = get_column_letter(col)
                html_parts.append(f'<th>{col_letter}</th>')
            html_parts.append('</tr></thead>')
            
            html_parts.append('<tbody>')
            
            for row_idx in range(1, max_row + 1):
                html_parts.append('<tr>')
                # Row header (1, 2, 3...)
                html_parts.append(f'<td class="row-header">{row_idx}</td>')
                
                for col_idx in range(1, max_col + 1):
                    cell_val = sheet_values.cell(row=row_idx, column=col_idx)
                    cell_formula = sheet_formulas.cell(row=row_idx, column=col_idx)
                    
                    value = cell_val.value
                    if value is None:
                        value = ''
                    
                    # Formula
                    formula = ""
                    if cell_formula.data_type == 'f':
                        formula = f"={cell_formula.value}"
                    elif value != "":
                        formula = str(value)
                        
                    # Styles
                    style_parts = []
                    
                    # Font styles
                    if cell_val.font:
                        if cell_val.font.bold:
                            style_parts.append('font-weight: bold')
                        if cell_val.font.italic:
                            style_parts.append('font-style: italic')
                        if cell_val.font.color and hasattr(cell_val.font.color, 'rgb') and cell_val.font.color.rgb:
                            # OpenPYXL sometimes returns '00000000' or similar
                            color = str(cell_val.font.color.rgb)
                            if len(color) == 8: # ARGB
                                color = '#' + color[2:]
                            elif len(color) == 6:
                                color = '#' + color
                            if color.startswith('#'):
                                style_parts.append(f'color: {color}')

                    # Background color
                    if cell_val.fill and isinstance(cell_val.fill, PatternFill) and cell_val.fill.start_color:
                        if hasattr(cell_val.fill.start_color, 'rgb') and cell_val.fill.start_color.rgb:
                            bg_color = str(cell_val.fill.start_color.rgb)
                            if len(bg_color) == 8: # ARGB
                                bg_color = '#' + bg_color[2:]
                            elif len(bg_color) == 6:
                                bg_color = '#' + bg_color
                            if bg_color.startswith('#') and bg_color != '#000000': # Ignore black default
                                style_parts.append(f'background-color: {bg_color}')
                    
                    style_attr = f'style="{"; ".join(style_parts)}"' if style_parts else ''
                    
                    html_parts.append(f'<td {style_attr} data-formula="{formula}" onclick="showFormula(this)">{value}</td>')
                
                html_parts.append('</tr>')
            
            html_parts.append('</tbody>')
            html_parts.append('</table>')
            html_parts.append('</div>')
            
            if sheet_values.max_row > max_rows:
                html_parts.append(f'<p class="text-muted small mt-2">Показано перші {max_rows} рядків з {sheet_values.max_row}</p>')
                
            html_parts.append('</div>') # End document-page
        
        html_content = '\n'.join(html_parts)
        return html_content, None

    except Exception as e:
        error_msg = f"Помилка при читанні таблиці: {str(e)}"
        return "", error_msg


def convert_pptx_to_html(file_path: str) -> Tuple[str, Optional[str]]:
    """
    Конвертує .pptx файл у HTML слайди
    
    Args:
        file_path: Шлях до .pptx файлу
        
    Returns:
        Tuple[str, Optional[str]]: (HTML контент, повідомлення про помилку)
    """
    try:
        from pptx import Presentation
        
        prs = Presentation(file_path)
        html_parts = []
        
        html_parts.append('<div class="pptx-viewer">')
        
        for slide_idx, slide in enumerate(prs.slides, 1):
            html_parts.append(f'<div class="slide-container card mb-4">')
            html_parts.append(f'<div class="card-header bg-primary text-white">')
            html_parts.append(f'<h5 class="mb-0">Слайд {slide_idx}</h5>')
            html_parts.append('</div>')
            html_parts.append('<div class="card-body">')
            
            # Витягти текст з усіх фігур на слайді
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text)
            
            if slide_texts:
                for text in slide_texts:
                    # Перший текст - заголовок
                    if text == slide_texts[0] and len(text) < 100:
                        html_parts.append(f'<h4>{text}</h4>')
                    else:
                        # Розбити на параграфи
                        paragraphs = text.split('\n')
                        for para in paragraphs:
                            if para.strip():
                                html_parts.append(f'<p>{para}</p>')
            else:
                html_parts.append('<p class="text-muted">Слайд не містить текстового контенту</p>')
            
            html_parts.append('</div>')
            html_parts.append('</div>')
        
        html_parts.append('</div>')
        
        html_content = '\n'.join(html_parts)
        
        if not prs.slides:
            html_content = '<p class="text-muted">Презентація не містить слайдів.</p>'
        
        return html_content, None
        
    except Exception as e:
        error_msg = f"Помилка при читанні презентації: {str(e)}"
        return "", error_msg

def convert_odt_to_html(file_path: str) -> Tuple[str, Optional[str]]:
    """
    Конвертує .odt файл у HTML
    """
    try:
        from odf.opendocument import load
        from odf.text import P, H
        from odf import teletype

        doc = load(file_path)
        html_parts = []

        # Обробка заголовків
        for h in doc.getElementsByType(H):
            text = teletype.extractText(h)
            level = h.getAttribute('outlinelevel') or 3
            html_parts.append(f'<h{level}>{text}</h{level}>')

        # Обробка параграфів
        for p in doc.getElementsByType(P):
            text = teletype.extractText(p)
            if text.strip():
                html_parts.append(f'<p>{text}</p>')

        html_content = '\n'.join(html_parts)
        
        if not html_content.strip():
            html_content = '<p class="text-muted">Документ порожній або не містить текстового контенту.</p>'
            
        return html_content, None

    except Exception as e:
        error_msg = f"Помилка при читанні ODT документа: {str(e)}"
        return "", error_msg


def convert_ods_to_html(file_path: str, max_rows: int = 100) -> Tuple[str, Optional[str]]:
    """
    Конвертує .ods файл у HTML таблиці
    """
    try:
        from odf.opendocument import load
        from odf.table import Table, TableRow, TableCell
        from odf import teletype

        doc = load(file_path)
        html_parts = []
        html_parts.append('<div class="document-page">')

        for sheet in doc.getElementsByType(Table):
            sheet_name = sheet.getAttribute('name')
            html_parts.append(f'<h4 class="mb-3">Аркуш: {sheet_name}</h4>')
            html_parts.append('<div class="excel-sheet mb-4">')
            html_parts.append('<table class="excel-table">')
            
            # Headers (A, B, C...) - ODS doesn't easily give us column count upfront, 
            # so we might skip the header row or try to estimate it. 
            # For simplicity, let's just render the data rows.
            
            rows = sheet.getElementsByType(TableRow)
            for i, row in enumerate(rows):
                if i >= max_rows:
                    break
                    
                html_parts.append('<tr>')
                # Add row number
                html_parts.append(f'<td class="row-header">{i+1}</td>')
                
                cells = row.getElementsByType(TableCell)
                for cell in cells:
                    # ODS uses 'number-columns-repeated' for empty cells or repeated values
                    repeat = int(cell.getAttribute('numbercolumnsrepeated') or 1)
                    
                    text = teletype.extractText(cell)
                    formula = cell.getAttribute('formula') or ""
                    
                    # Clean up formula (ODS formulas often start with 'of:=')
                    if formula.startswith('of:='):
                        formula = '=' + formula[4:]
                    elif not formula and text:
                        formula = text
                        
                    for _ in range(repeat):
                        html_parts.append(f'<td data-formula="{formula}" onclick="showFormula(this)">{text}</td>')
                        
                html_parts.append('</tr>')
            
            html_parts.append('</table>')
            html_parts.append('</div>')
            
            if len(rows) > max_rows:
                html_parts.append(f'<p class="text-muted small">Показано перші {max_rows} рядків з {len(rows)}</p>')

        html_parts.append('</div>')
        
        html_content = '\n'.join(html_parts)
        return html_content, None

    except Exception as e:
        error_msg = f"Помилка при читанні ODS таблиці: {str(e)}"
        return "", error_msg


def convert_odp_to_html(file_path: str) -> Tuple[str, Optional[str]]:
    """
    Конвертує .odp файл у HTML (витягує текст зі слайдів)
    """
    try:
        from odf.opendocument import load
        from odf.draw import Page, Frame, TextBox
        from odf import teletype

        doc = load(file_path)
        html_parts = []
        
        html_parts.append('<div class="pptx-viewer">')

        slides = doc.getElementsByType(Page)
        for i, slide in enumerate(slides, 1):
            html_parts.append(f'<div class="slide-container card mb-4">')
            html_parts.append(f'<div class="card-header bg-primary text-white">')
            html_parts.append(f'<h5 class="mb-0">Слайд {i}</h5>')
            html_parts.append('</div>')
            html_parts.append('<div class="card-body">')
            
            # Extract text from frames/textboxes
            slide_text_found = False
            for frame in slide.getElementsByType(Frame):
                for textbox in frame.getElementsByType(TextBox):
                    text = teletype.extractText(textbox)
                    if text.strip():
                        html_parts.append(f'<p>{text}</p>')
                        slide_text_found = True
            
            if not slide_text_found:
                html_parts.append('<p class="text-muted">Слайд не містить текстового контенту</p>')
                
            html_parts.append('</div>')
            html_parts.append('</div>')

        html_parts.append('</div>')
        html_content = '\n'.join(html_parts)
        
        if not slides:
             html_content = '<p class="text-muted">Презентація не містить слайдів.</p>'

        return html_content, None

    except Exception as e:
        error_msg = f"Помилка при читанні ODP презентації: {str(e)}"
        return "", error_msg


def get_archive_content(file_path: str, file_ext: str) -> Tuple[list, Optional[str]]:
    """
    Повертає список файлів в архіві
    """
    files_list = []
    try:
        if file_ext == '.zip':
            import zipfile
            with zipfile.ZipFile(file_path, 'r') as zf:
                for info in zf.infolist():
                    files_list.append({
                        'name': info.filename,
                        'size': info.file_size,
                        'is_dir': info.is_dir()
                    })
                    
        elif file_ext == '.rar':
            import rarfile
            with rarfile.RarFile(file_path) as rf:
                for info in rf.infolist():
                    files_list.append({
                        'name': info.filename,
                        'size': info.file_size,
                        'is_dir': info.isdir()
                    })
                    
        elif file_ext == '.7z':
            import py7zr
            with py7zr.SevenZipFile(file_path, 'r') as zf:
                for info in zf.list():
                    files_list.append({
                        'name': info.filename,
                        'size': info.uncompressed,
                        'is_dir': info.is_directory
                    })
                    
        return files_list, None
        
    except Exception as e:
        error_msg = f"Помилка при читанні архіву: {str(e)}"
        return [], error_msg

def get_file_type_info(file_ext: str) -> dict:
    """
    Повертає інформацію про тип файлу базуючись на розширенні
    
    Args:
        file_ext: Розширення файлу (наприклад, '.docx')
        
    Returns:
        dict: Словник з інформацією про тип файлу
    """
    file_types = {
        # Office документи
        '.doc': {'name': 'Word документ', 'icon': '📝', 'color': 'soft-blue', 'preview': False},
        '.docx': {'name': 'Word документ', 'icon': '📝', 'color': 'soft-blue', 'preview': True},
        '.odt': {'name': 'OpenDocument Текст', 'icon': '📝', 'color': 'soft-blue', 'preview': True},
        
        '.xls': {'name': 'Excel таблиця', 'icon': '📊', 'color': 'soft-green', 'preview': False},
        '.xlsx': {'name': 'Excel таблиця', 'icon': '📊', 'color': 'soft-green', 'preview': True},
        '.ods': {'name': 'OpenDocument Таблиця', 'icon': '📊', 'color': 'soft-green', 'preview': True},
        
        '.ppt': {'name': 'PowerPoint', 'icon': '📽️', 'color': 'soft-orange', 'preview': False},
        '.pptx': {'name': 'PowerPoint', 'icon': '📽️', 'color': 'soft-orange', 'preview': True},
        '.odp': {'name': 'OpenDocument Презентація', 'icon': '📽️', 'color': 'soft-orange', 'preview': True},
        
        # PDF
        '.pdf': {'name': 'PDF документ', 'icon': '📄', 'color': 'soft-red', 'preview': True},
        
        # Зображення
        '.jpg': {'name': 'Зображення', 'icon': '🖼️', 'color': 'soft-purple', 'preview': True},
        '.jpeg': {'name': 'Зображення', 'icon': '🖼️', 'color': 'soft-purple', 'preview': True},
        '.png': {'name': 'Зображення', 'icon': '🖼️', 'color': 'soft-purple', 'preview': True},
        '.gif': {'name': 'Зображення', 'icon': '🖼️', 'color': 'soft-purple', 'preview': True},
        '.bmp': {'name': 'Зображення', 'icon': '🖼️', 'color': 'soft-purple', 'preview': True},
        '.webp': {'name': 'Зображення', 'icon': '🖼️', 'color': 'soft-purple', 'preview': True},
        
        # Текстові файли
        '.txt': {'name': 'Текстовий файл', 'icon': '📃', 'color': 'soft-gray', 'preview': True},
        
        # Код
        '.py': {'name': 'Python', 'icon': '🐍', 'color': 'soft-teal', 'preview': True},
        '.html': {'name': 'HTML', 'icon': '💻', 'color': 'soft-teal', 'preview': True},
        '.css': {'name': 'CSS', 'icon': '💻', 'color': 'soft-teal', 'preview': True},
        '.js': {'name': 'JavaScript', 'icon': '💻', 'color': 'soft-teal', 'preview': True},
        '.json': {'name': 'JSON', 'icon': '💻', 'color': 'soft-teal', 'preview': True},
        '.xml': {'name': 'XML', 'icon': '💻', 'color': 'soft-teal', 'preview': True},
        
        # Архіви
        '.zip': {'name': 'ZIP архів', 'icon': '📦', 'color': 'soft-gray', 'preview': True},
        '.rar': {'name': 'RAR архів', 'icon': '📦', 'color': 'soft-gray', 'preview': True},
        '.7z': {'name': '7Z архів', 'icon': '📦', 'color': 'soft-gray', 'preview': True},
    }
    
    return file_types.get(file_ext.lower(), {
        'name': f'{file_ext.upper()} файл',
        'icon': '📎',
        'color': 'soft-gray',
        'preview': False
    })
